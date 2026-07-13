"""font-aware 横幅計測の回帰テスト。

`validate_fit` の latin 幅計測が「要素ごとの実フォント」を参照することを確認する
（旧実装は family を無視し常に Arial で計測していた）。

実フォントのメトリクスに依存せず決定的に回すため、幅判定を左右するコアの
`_latin_w` をスタブ化し、family ごとに異なる字幅を与えて
レジストリ →（font 伝播）→ `validate_fit` の判定変化を検証する。
実フォント解決（fc-match）は環境依存なので、その差分テストは skip ガードで囲む。
"""
import pytest
from pptx.enum.shapes import MSO_SHAPE

import slides
from slides import shape_box, textbox, validate_fit, PT_PER_CM


def _overflow_h(prs):
    return [v for v in validate_fit(prs) if v.code == "OVERFLOW_H"]


def test_register_records_element_font(prs, body_slide):
    """ビルダーの font= がレジストリに記録される。"""
    slides._FIT_REGISTRY.clear()
    textbox(body_slide, 1.07, 5.0, 5.0, 0.8, "X", size=10, font="Futura")
    assert slides._FIT_REGISTRY[-1]["font"] == "Futura"


def test_register_defaults_to_latin_FONT(prs, body_slide):
    """font 未指定は既定 latin フォント FONT で記録される（後方互換）。"""
    slides._FIT_REGISTRY.clear()
    textbox(body_slide, 1.07, 5.0, 5.0, 0.8, "X", size=10)
    assert slides._FIT_REGISTRY[-1]["font"] == slides.FONT


def test_validate_fit_uses_element_font(prs, monkeypatch):
    """同じ文字列・同じ箱でも、要素 font で横はみ出し判定が変わる。

    実フォント非依存にするため `_latin_w` をスタブ化：
      family=="Wide" は 12pt/char、それ以外は 6pt/char。
    箱の有効幅を Narrow がちょうど収まる幅に設定し、
      - font=Narrow → OVERFLOW_H なし
      - font=Wide   → OVERFLOW_H あり（旧実装なら Arial 計測で見逃していた）
    を確認する。
    """
    def fake_latin_w(s, size_pt, bold, family=None):
        per = 12.0 if family == "Wide" else 6.0
        return len(s) * per

    monkeypatch.setattr(slides, "_latin_w", fake_latin_w)

    token = "ABCDE"                     # 5 文字
    narrow_pt = len(token) * 6.0        # Narrow の実幅 = 30pt
    box_w_cm = narrow_pt / PT_PER_CM + 0.24   # 有効幅=Narrow実幅（margin 0.12x2）

    def build(font):
        slides._FIT_REGISTRY.clear()
        s = prs.slides.add_slide(prs.slide_layouts[slides.L_BODY])
        shape_box(s, MSO_SHAPE.RECTANGLE, 3.0, 5.0, box_w_cm, 0.7,
                  text=token, size=10, bold=True, fill=None, font=font)
        return _overflow_h(prs)

    assert not build("Narrow"), "Narrow は箱に収まるべき"
    assert build("Wide"), "Wide は横はみ出しを検知すべき（font 伝播）"


def test_resolve_arial_returns_existing_file():
    """Arial 解決は実在ファイルを返す（既定ブランドの計測が壊れていない）。
    fontconfig も同梱 Arial も無い極小環境では skip。"""
    path, _matched = slides._resolve_font_path("Arial", False)
    if path is None:
        pytest.skip("フォント解決系（fontconfig / 同梱 Arial）が無い環境")
    import os
    assert os.path.exists(path)


def test_real_font_widths_differ_when_available():
    """区別できる実フォントが 2 種見える環境では、_latin_w は family で
    異なる幅を返す（font-aware であることの実測確認）。無ければ skip。"""
    pa, _ = slides._resolve_font_path("Arial", True)
    pv, _ = slides._resolve_font_path("Verdana", True)
    if not pa or not pv or pa == pv:
        pytest.skip("区別できる実フォントが揃わない環境")
    wa = slides._latin_w("RECOMMENDED", 12, True, "Arial")
    wv = slides._latin_w("RECOMMENDED", 12, True, "Verdana")
    assert wa != wv
