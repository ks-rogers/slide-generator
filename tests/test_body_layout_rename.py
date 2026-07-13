"""本文レイアウト判定の脱・名前依存の回帰テスト。

ゾーン検査（ZONE）の発火条件はかつてレイアウト名リテラル
（"タイトル社外秘"）への部分一致で、リブランドで本文レイアウトを
改名すると検査ブロックごと無音でスキップされ、画面外はみ出しでも
PASS してしまった。判定は prs.slide_layouts[L_BODY].name からの
動的解決（REBRAND.md が契約とする index 構成）に変更済みで、
レイアウト名が何であれ検知が生きることを保証する。
"""
from pptx.oxml.ns import qn

from slides import (
    add_page_number, connector, textbox, validate_fit, L_BODY, WHITE, PRIMARY,
)


def _rename_body_layout(prs, name):
    """リブランドで本文レイアウトを改名した状況を再現する。"""
    layout = prs.slide_layouts[L_BODY]
    layout._element.find(qn("p:cSld")).set("name", name)


def _zone_violations(prs):
    return [v for v in validate_fit(prs) if v.code == "ZONE"]


def _add_out_of_zone_textbox(slide):
    # 下端 13.8 + 0.6 = 14.4cm > Y_MAX(12.8) → ZONE ERROR になるべき配置
    textbox(slide, 1.07, 13.8, 20.0, 0.6, "ゾーン外の文字", size=11)


# ------------------------------------------------------- _FIT_REGISTRY (textbox)
def test_zone_detected_with_default_layout_name(prs, body_slide):
    _add_out_of_zone_textbox(body_slide)
    assert len(_zone_violations(prs)) == 1


def test_zone_detected_after_layout_rename(prs, body_slide):
    """本文レイアウトを改名しても ZONE 検知が無音化しない。"""
    _rename_body_layout(prs, "Body")
    _add_out_of_zone_textbox(body_slide)
    assert len(_zone_violations(prs)) == 1


# ------------------------------------------------------ _GEOM_REGISTRY (connector)
def test_geom_zone_detected_after_layout_rename(prs, body_slide):
    """非テキスト要素（コネクタ等）のゾーン判定も改名に追従する。"""
    _rename_body_layout(prs, "Body")
    connector(body_slide, 20.0, 4.0, 30.0, 4.0)  # 右端 30cm > X_MAX(24.33)
    assert any(v.kind == "connector" for v in _zone_violations(prs))


# ------------------------------------------------------------- add_page_number
def test_page_number_color_follows_rename(prs, body_slide):
    """本文スライドのページ番号色（WHITE）判定も改名に追従する。"""
    _rename_body_layout(prs, "Body")
    tb = add_page_number(body_slide, 2)
    run = tb.text_frame.paragraphs[0].runs[0]
    assert run.font.color.rgb == WHITE


def test_page_number_color_on_chapter_layout(prs):
    """章扉（非本文）は accent 色のまま（判定が過剰マッチしない）。"""
    from slides import L_CHAPTER
    chapter = prs.slides.add_slide(prs.slide_layouts[L_CHAPTER])
    tb = add_page_number(chapter, 2)
    run = tb.text_frame.paragraphs[0].runs[0]
    assert run.font.color.rgb == PRIMARY
