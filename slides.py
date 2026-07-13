"""slides.py — スライドテンプレート用 共通ヘルパー

雛形（templates/ 配下、ファイル名は brand.py の TEMPLATE）に依存する再利用可能な
ビルディングブロック。各案件の generate.py からこのモジュールを import して
ページ定義のみ書く運用を想定。

主な要素:
  - ブランドカラー定数（PRIMARY/SECONDARY/SUCCESS/DANGER/HIGHLIGHT/TEXT/TEXT_MUTED/BORDER/SURFACE）
  - レイアウト定数（L_BACK/L_CHAPTER/L_BODY/L_COVER）
  - 低レベル: textbox / shape_box / connector / picture / fill_tf
  - グラフ: bar_chart / line_chart / pie_chart（ブランド配色）
  - スライドタイプ設定: configure_body / configure_chapter / update_cover
  - その他: add_page_number / delete_slide

スライドサイズ: 25.40 x 14.29 cm (16:9)
本文有効エリア: x 1.07–24.33, y 3.0–12.8 (CONFIDENTIAL の 12.92cm 上まで)
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree as _etree

# ---------------------------------------------------------------- brand
# 色・フォント・既定雛形ファイル名は brand.py に一元化（リブランドの単一ソース）。
# 他モジュール（generate.py 等）は従来どおり `from slides import PRIMARY, ...` でよい
# ＝ここで再エクスポートする。リブランドは brand.py を編集する（docs/REBRAND.md）。
from brand import (
    PRIMARY, PRIMARY_LIGHT, SECONDARY, SUCCESS, DANGER, HIGHLIGHT,
    TEXT_MUTED, TEXT, BORDER, SURFACE, WHITE, BLACK,
    FONT, JP_FONT, JOSEFIN, TEMPLATE,
)

# タイプスケール（デザインシステム px → PPTX pt 換算 ×0.75）。本文コンポーネント推奨値。
T_H1 = 24.0       # ページ大タイトル（chrome 側。参考値）
T_H2 = 15.0       # カード見出し等
T_BODY = 10.5     # 本文
T_CAPTION = 8.25  # キャプション・補足

# ---------------------------------------------------------------- canvas
SLIDE_W_CM = 25.40
SLIDE_H_CM = 14.29

# ---------------------------------------------------------------- layouts (v3)
L_BACK = 0      # 裏紙（裏表紙）
L_CHAPTER = 1   # TITLE_AND_BODY: 章扉（全面主色）
L_BODY = 2      # ページ番号+タイトル社外秘: 通常本文
L_COVER = 3     # TITLE: 表紙

# ---------------------------------------------------------------- xml ns
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# ---------------------------------------------------------------- content zone
X_MIN, X_MAX = 1.07, 24.33     # 本文コンテンツ有効 x（cm）
Y_MIN, Y_MAX = 3.0, 12.8       # 本文コンテンツ有効 y（cm）
PT_PER_CM = 28.3465            # 1cm = 28.3465pt（PDF/pptx 共通）

# ============================================================================
# レイアウト検知レジストリ
#   各ビルダーが「描いたテキスト要素」を記録する。レンダリング前後で
#   枠はみ出し・見切れ・ゾーン逸脱を機械的に検査するために使う。
#   描画は一切変更しない（記録の追記のみ）。
# ============================================================================

_FIT_REGISTRY: list = []

# 塗りつぶしの矩形「帯」（テキストなしの色帯ヘッダ・バッジ等）を記録する。
# テキスト要素ではないため _FIT_REGISTRY には載らないが、「textbox が
# 色帯の下端を突き抜けて漏出していないか」を判定するための背景として使う。
_BAND_REGISTRY: list = []

# 非テキストの幾何要素（コネクタ線・画像・グラフ）の bbox を記録する。
# テキストを持たないため _FIT_REGISTRY には載らないが、本文コンテンツゾーン
# （x 1.07–24.33 / y 3.0–12.8）からの逸脱を validate() で検査するために使う。
_GEOM_REGISTRY: list = []

# 矢じり付きコネクタ（矢印）の軸長・線幅を記録する。矢印は線（geom）なので
# _FIT_REGISTRY には載らず、「バンド間の隙間が狭すぎて軸が無く矢じりだけ潰れる
# ＝向きが読めない矢印」は従来の検査では拾えなかった。validate_fit が「軸の
# 可視長」不足を WARN(ARROW_TOO_SHORT) で報告するために使う。
_ARROW_REGISTRY: list = []

# カード/パネル等の「背景枠」を register_container() で親として登録する。
# validate_fit が「枠内に置いたはずの子要素（チップ・テキスト・図形）が枠外へ
# はみ出していないか」を WARN(CONTAINER_OVERFLOW) で検知するために使う。
# geom のゾーン検査（ページ枠）では拾えない「枠内のはみ出し」を補完する。
_CONTAINER_REGISTRY: list = []


def _register(slide, kind, x, y, w, h, lines, *, size, bold=False,
              line_spacing=None, ml=0.1, mr=0.1, mt=0.05, mb=0.05,
              autosize_v=False, font=None):
    """テキスト要素をレジストリに記録する。失敗しても生成は止めない
    （ただし沈黙はせず stderr に警告を出す）。

    kind: 'textbox'（縦に自動伸長）/ 'shape'（クリップ）/
          'cell'（表セル・クリップ）/ 'ph_fixed'（テンプレ固定枠・ゾーン免除）
    autosize_v: True ならビルダー側が枠高を内容に合わせて自動調整するため
                生成時の縦あふれ検査を行わない（横/見切れ実測は継続）。
    """
    try:
        if isinstance(lines, str):
            lines = [lines]
        lines = [str(s) for s in lines if str(s) != ""]
        if not lines:
            return
        _FIT_REGISTRY.append({
            "slide_id": slide.slide_id,
            "layout": slide.slide_layout.name,
            "kind": kind,
            "x": float(x), "y": float(y), "w": float(w), "h": float(h),
            "lines": lines,
            "size": float(size),
            "bold": bool(bold),
            # latin 幅計測に使う実フォント（要素ごとの font= 上書きを反映）。
            # 未指定は既定 latin フォント FONT。
            "font": str(font) if font else FONT,
            "ls": float(line_spacing) if line_spacing else 1.0,
            "ml": float(ml), "mr": float(mr),
            "mt": float(mt), "mb": float(mb),
            "autosize_v": bool(autosize_v),
        })
    except Exception as exc:  # noqa: BLE001  記録失敗を握り潰さない
        import sys
        print(f"[slides._register] 記録失敗 ({kind}): {exc}",
              file=sys.stderr)


def _register_band(slide, x, y, w, h):
    """塗りつぶし矩形（色帯）の矩形を記録する。失敗しても生成は止めない。"""
    try:
        _BAND_REGISTRY.append({
            "slide_id": slide.slide_id,
            "x": float(x), "y": float(y), "w": float(w), "h": float(h),
        })
    except Exception as exc:  # noqa: BLE001  記録失敗を握り潰さない
        import sys
        print(f"[slides._register_band] 記録失敗: {exc}", file=sys.stderr)


def _register_geom(slide, kind, x, y, w, h):
    """非テキストの幾何要素（コネクタ/画像/グラフ）の bbox を記録する。
    失敗しても生成は止めない（沈黙はせず stderr に警告を出す）。

    kind: 'connector' / 'picture' / 'chart'。
    """
    try:
        _GEOM_REGISTRY.append({
            "slide_id": slide.slide_id,
            "layout": slide.slide_layout.name,
            "kind": kind,
            "x": float(x), "y": float(y), "w": float(w), "h": float(h),
        })
    except Exception as exc:  # noqa: BLE001  記録失敗を握り潰さない
        import sys
        print(f"[slides._register_geom] 記録失敗 ({kind}): {exc}",
              file=sys.stderr)


def _register_arrow(slide, x, y, w, h, length_cm, width_pt, n_heads):
    """矢じり付きコネクタ（矢印）の軸長・線幅・矢じり数を記録する。
    失敗しても生成は止めない（沈黙はせず stderr に警告を出す）。

    validate_fit が「総長 − 矢じり長 = 軸の可視長」を求め、軸が短すぎて
    矢じりだけが潰れる矢印を ARROW_TOO_SHORT WARN として検知するために使う。
    """
    try:
        _ARROW_REGISTRY.append({
            "slide_id": slide.slide_id,
            "x": float(x), "y": float(y), "w": float(w), "h": float(h),
            "length": float(length_cm),
            "width": float(width_pt),
            "heads": int(n_heads),
        })
    except Exception as exc:  # noqa: BLE001  記録失敗を握り潰さない
        import sys
        print(f"[slides._register_arrow] 記録失敗: {exc}", file=sys.stderr)


def register_container(slide, x, y, w, h):
    """カード/パネル等の「背景枠」を子要素はみ出し検査の親として登録する。

    背景矩形（``card()`` や各案件の panel ヘルパー等）を登録しておくと、
    validate_fit が「中心がこの枠内にある子要素（チップ・テキスト・図形）」を
    集め、子の bbox が枠を越えていれば CONTAINER_OVERFLOW(WARN) を報告する。
    描画は一切しない（記録のみ）。失敗しても生成は止めない。

    Args:
        slide: 対象スライド。
        x, y, w, h: 背景枠の矩形（cm）。
    """
    try:
        _CONTAINER_REGISTRY.append({
            "slide_id": slide.slide_id,
            "x": float(x), "y": float(y), "w": float(w), "h": float(h),
        })
    except Exception as exc:  # noqa: BLE001  記録失敗を握り潰さない
        import sys
        print(f"[slides.register_container] 記録失敗: {exc}", file=sys.stderr)


# ============================================================================
# low-level utilities
# ============================================================================

def delete_slide(prs, idx):
    """index指定でスライドを削除（rels もクリーンアップ）"""
    sld_lst = prs.slides._sldIdLst
    sld = list(sld_lst)[idx]
    rid = sld.get(R_NS + "id")
    sld_lst.remove(sld)
    prs.part.drop_rel(rid)


def _set_run(run, text, *, size=10, bold=False, color=TEXT, font=FONT,
             italic=False, jp_font=JP_FONT):
    """Run（文字列セグメント）にテキストとスタイルを設定する。

    latin（英数字）は ``font``、和文（ea/cs）は ``jp_font`` を割り当てる。こうする
    ことで英字は Arial、和文は Noto Sans JP と明示制御でき、Arial の無制御 CJK
    フォールバック（MS P ゴシック等）に落ちる古い見た目を防ぐ。"""
    run.text = text
    run.font.name = font          # <a:latin>
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # 和文（east-asian）・complex-script は jp_font を割り当てる
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = _etree.SubElement(rPr, qn(tag))
        el.set('typeface', jp_font)


def fill_tf(tf, lines, *, size=10, bold=False, color=TEXT,
            align=PP_ALIGN.LEFT, line_spacing=None, font=FONT, anchor=None,
            jp_font=JP_FONT):
    """テキストフレームを多行テキストで埋める。各行は別段落。"""
    if not isinstance(lines, list):
        lines = [lines]
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    paras = list(tf.paragraphs)
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    p0 = paras[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    p0.alignment = align
    if line_spacing:
        p0.line_spacing = line_spacing
    run = p0.add_run()
    _set_run(run, lines[0], size=size, bold=bold, color=color, font=font, jp_font=jp_font)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        _set_run(run, line, size=size, bold=bold, color=color, font=font, jp_font=jp_font)


# ============================================================================
# building blocks
# ============================================================================

def textbox(slide, x, y, w, h, text, *, size=10, bold=False, color=TEXT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None,
            border=None, border_w=0.5, line_spacing=None, italic=False, font=FONT,
            jp_font=JP_FONT):
    """テキストボックスをスライドに追加する。

    Args:
        slide: 追加先スライドオブジェクト。
        x, y: 左上座標（cm）。
        w, h: 幅・高さ（cm）。縦方向は内容に合わせて自動伸長する。
        text: 文字列または文字列リスト（複数行）。
        size: フォントサイズ（pt）。
        bold: 太字フラグ。
        color: 文字色（RGBColor）。
        align: 横揃え（PP_ALIGN.*）。
        anchor: 縦揃え（MSO_ANCHOR.*）。
        fill: 背景塗りつぶし色（RGBColor）。None で透明。
        border: 枠線色（RGBColor）。None で非表示。
        border_w: 枠線幅（pt）。
        line_spacing: 行間倍率（float）。None でデフォルト。
        italic: イタリック体フラグ。
        font: フォント名。

    Returns:
        追加された Shape オブジェクト。
    """
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(border_w)
    else:
        tb.line.fill.background()
    # テキストボックスにテーマ由来の影を継承させない（文字へのドロップシャドウ防止）
    tb.shadow.inherit = False
    fill_tf(tf, text, size=size, bold=bold, color=color,
            align=align, line_spacing=line_spacing, anchor=anchor, font=font, jp_font=jp_font)
    if italic:
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.italic = True
    _register(slide, "textbox", x, y, w, h, text, size=size, bold=bold,
              line_spacing=line_spacing, ml=0.1, mr=0.1, mt=0.05, mb=0.05,
              font=font)
    return tb


def _kill_shadow(sh):
    """シェイプからテーマ継承のドロップシャドウを完全に除去する。

    ``shadow.inherit=False``（spPr に空 effectLst を挿入）だけでは
    プリセット図形の ``<p:style><a:effectRef>`` 経由のテーマ影が
    LibreOffice/PowerPoint で残ることがあるため、effectRef の idx も 0 にする。"""
    sh.shadow.inherit = False
    style = sh._element.find(qn("p:style"))
    if style is not None:
        eff = style.find(qn("a:effectRef"))
        if eff is not None:
            eff.set("idx", "0")
    return sh


def shape_box(slide, shape_type, x, y, w, h, *, text=None, fill=None,
              line=None, line_w=0.5, size=10, bold=False, color=TEXT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=None, font=FONT, jp_font=JP_FONT, shadow=False):
    """任意の MSO_SHAPE 図形をスライドに追加する。テキスト付き可。

    Args:
        slide: 追加先スライドオブジェクト。
        shape_type: 図形種別（MSO_SHAPE.*）。
        x, y: 左上座標（cm）。
        w, h: 幅・高さ（cm）。固定サイズ（テキストははみ出しクリップ）。
        text: 文字列または文字列リスト。None でテキストなし。
        fill: 塗りつぶし色（RGBColor）。None で透明。
        line: 枠線色（RGBColor）。None で非表示。
        line_w: 枠線幅（pt）。
        size: フォントサイズ（pt）。
        bold: 太字フラグ。
        color: 文字色（RGBColor）。
        align: 横揃え（PP_ALIGN.*）。
        anchor: 縦揃え（MSO_ANCHOR.*）。
        line_spacing: 行間倍率（float）。
        font: latin（英数字）フォント名。
        jp_font: 和文フォント名。
        shadow: ドロップシャドウを残すか。既定 False（フラット）。True で
            塗りつぶしカードにテーマ影を付与（従来挙動の opt-in）。

    Returns:
        追加された Shape オブジェクト。
    """
    sh = slide.shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    # フラット標準: 既定で影を除去する。shadow=True かつ塗りつぶしありの場合のみ
    # テーマのドロップシャドウを残す（カード影を明示的に欲しいケースの opt-in）。
    # 透明シェイプ（fill=None）は文字グリフに影が落ちるため常に除去する。
    if not (shadow and fill is not None):
        _kill_shadow(sh)
    # 塗りつぶし矩形は「帯」候補として記録（textbox の帯下端漏出検査の背景）。
    if fill is not None and shape_type == MSO_SHAPE.RECTANGLE:
        _register_band(slide, x, y, w, h)
    if text is not None:
        tf = sh.text_frame
        tf.margin_left = Cm(0.12)
        tf.margin_right = Cm(0.12)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)
        fill_tf(tf, text, size=size, bold=bold, color=color,
                align=align, line_spacing=line_spacing, anchor=anchor, font=font, jp_font=jp_font)
        _register(slide, "shape", x, y, w, h, text, size=size, bold=bold,
                  line_spacing=line_spacing, ml=0.12, mr=0.12, mt=0.05, mb=0.05,
                  font=font)
    return sh


def connector(slide, x1, y1, x2, y2, *, color=TEXT_MUTED, width=0.75,
              dash=None, connector_type=MSO_CONNECTOR.STRAIGHT,
              begin_arrow=False, end_arrow=False, shadow=False):
    """直線/コネクタをスライドに追加する（cm 座標）。

    概念図・フロー図で「ボックス同士を線でつなぐ」用途。テキストは持たない。
    始点 (x1,y1)・終点 (x2,y2) を囲む矩形を geom レジストリに記録し、本文
    コンテンツゾーンからの逸脱を validate() で検査できるようにする。

    Args:
        slide: 追加先スライド。
        x1, y1: 始点座標（cm）。
        x2, y2: 終点座標（cm）。
        color: 線色（RGBColor）。ブランド色のみ。既定 TEXT_MUTED。
        width: 線幅（pt）。
        dash: 破線スタイル（MSO_LINE_DASH_STYLE.* / None で実線）。
        connector_type: MSO_CONNECTOR.*（STRAIGHT / ELBOW / CURVE）。
        begin_arrow: 始点に矢じりを付ける。
        end_arrow: 終点に矢じりを付ける。
        shadow: ドロップシャドウを残すか。既定 False（フラット＝影なし）。
            True のときのみテーマの影を継承する。

    Returns:
        追加された Connector オブジェクト。
    """
    sh = slide.shapes.add_connector(
        connector_type, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    sh.line.color.rgb = color
    sh.line.width = Pt(width)
    if dash is not None:
        sh.line.dash_style = dash
    # 線・矢印はフラット標準で影なし。spPr の空 effectLst（shadow.inherit=False）
    # だけではコネクタの <p:style><a:effectRef> 経由のテーマ影が LibreOffice /
    # PowerPoint で残るため、_kill_shadow で effectRef の idx も 0 にして除去する。
    # どうしても立体感が要る場合のみ shadow=True でテーマ影をそのまま継承する。
    if not shadow:
        _kill_shadow(sh)
    if begin_arrow or end_arrow:
        ln = sh.line._get_or_add_ln()
        # スキーマ上 headEnd → tailEnd の順序で末尾に追加する
        if begin_arrow:
            he = _etree.SubElement(ln, qn('a:headEnd'))
            he.set('type', 'triangle'); he.set('w', 'med'); he.set('len', 'med')
        if end_arrow:
            te = _etree.SubElement(ln, qn('a:tailEnd'))
            te.set('type', 'triangle'); te.set('w', 'med'); te.set('len', 'med')
    _register_geom(slide, "connector",
                   min(x1, x2), min(y1, y2), abs(x2 - x1), abs(y2 - y1))
    if begin_arrow or end_arrow:
        dx, dy = x2 - x1, y2 - y1
        length = (dx * dx + dy * dy) ** 0.5
        _register_arrow(slide, min(x1, x2), min(y1, y2), abs(dx), abs(dy),
                        length, width, int(begin_arrow) + int(end_arrow))
    return sh


def picture(slide, image, x, y, w=None, h=None):
    """画像（写真・ロゴ・図版）をスライドに配置する（cm 座標）。

    w/h は省略可。片方だけ指定するとアスペクト比を保って他方を自動算出する
    （両方 None なら画像の原寸）。配置後の実寸を geom レジストリに記録し、
    本文コンテンツゾーン逸脱を validate() で検査できるようにする。

    Args:
        slide: 追加先スライド。
        image: 画像ファイルのパス（str / Path）またはファイルライクオブジェクト。
        x, y: 左上座標（cm）。
        w: 幅（cm）。None で自動（h かアスペクト比から決定）。
        h: 高さ（cm）。None で自動。

    Returns:
        追加された Picture オブジェクト。

    Raises:
        FileNotFoundError: image がパスで、ファイルが存在しない場合。
    """
    if isinstance(image, (str, Path)):
        p = Path(image)
        if not p.exists():
            raise FileNotFoundError(f"画像が見つかりません: {p}")
        image = str(p)
    pic = slide.shapes.add_picture(
        image, Cm(x), Cm(y),
        Cm(w) if w is not None else None,
        Cm(h) if h is not None else None)
    _register_geom(slide, "picture",
                   Emu(pic.left).cm, Emu(pic.top).cm,
                   Emu(pic.width).cm, Emu(pic.height).cm)
    return pic


# ============================================================================
# chart helpers  (棒 / 折れ線 / 円。ブランド配色で整形)
#   ※ グラフ内部テキスト（軸・データラベル）は validate_render の文字照合対象外。
#     配置（コンテンツゾーン）は _register_geom 経由で validate_fit が検査する。
# ============================================================================

# ブランド配色のグラフ用既定パレット（系列／扇の順に適用）
CHART_PALETTE = (PRIMARY, SECONDARY, SUCCESS, HIGHLIGHT, TEXT_MUTED, DANGER)


def _add_chart(slide, chart_type, x, y, w, h, chart_data):
    """チャートを追加し geom レジストリに bbox を記録、GraphicFrame を返す（内部用）。"""
    gf = slide.shapes.add_chart(
        chart_type, Cm(x), Cm(y), Cm(w), Cm(h), chart_data)
    _register_geom(slide, "chart", x, y, w, h)
    return gf


def _style_axis_fonts(chart, *, font_size, color=TEXT):
    """カテゴリ軸・値軸の目盛ラベルフォントをブランド標準（Arial / TEXT）に整える。"""
    for axis in (chart.category_axis, chart.value_axis):
        font = axis.tick_labels.font
        font.size = Pt(font_size)
        font.name = FONT
        font.color.rgb = color


def _style_legend(chart, *, position, font_size, color=TEXT):
    """凡例をブランド標準で表示する。"""
    chart.has_legend = True
    chart.legend.position = position
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(font_size)
    chart.legend.font.name = FONT
    chart.legend.font.color.rgb = color


def _style_data_labels(plot_or_dl, *, font_size, number_format=None, color=TEXT):
    """データラベルのフォント・書式をブランド標準に整える。"""
    dl = plot_or_dl
    dl.font.size = Pt(font_size)
    dl.font.name = FONT
    dl.font.color.rgb = color
    if number_format:
        dl.number_format = number_format
        dl.number_format_is_linked = False


def _build_category_data(categories, series, number_format):
    """CategoryChartData を組み立てる。series は (系列名, 値リスト) のリスト。"""
    cd = CategoryChartData()
    cd.categories = categories
    for name, values in series:
        if number_format:
            cd.add_series(name, values, number_format=number_format)
        else:
            cd.add_series(name, values)
    return cd


def bar_chart(slide, x, y, w, h, *, categories, series, colors=None,
              horizontal=False, show_value=True, show_legend=None,
              number_format=None, font_size=10):
    """棒グラフ（縦／横）を配置する。

    Args:
        slide: 追加先スライド。
        x, y, w, h: 配置矩形（cm）。
        categories: カテゴリ名のリスト（軸ラベル）。
        series: (系列名, 値リスト) のタプルのリスト。
                例: [("2024", [10, 20, 30]), ("2025", [15, 25, 35])]
        colors: 系列ごとの色（RGBColor）のリスト。省略時は CHART_PALETTE。
        horizontal: True で横棒（BAR_CLUSTERED）、False で縦棒（COLUMN_CLUSTERED）。
        show_value: データラベル（値）を表示する。
        show_legend: 凡例表示。None なら系列数>1 のとき自動で表示。
        number_format: 値の表示書式（例 '#,##0' / '0.0%'）。
        font_size: 軸・ラベル・凡例の基準フォントサイズ（pt）。

    Returns:
        Chart オブジェクト。
    """
    ct = (XL_CHART_TYPE.BAR_CLUSTERED if horizontal
          else XL_CHART_TYPE.COLUMN_CLUSTERED)
    cd = _build_category_data(categories, series, number_format)
    chart = _add_chart(slide, ct, x, y, w, h, cd).chart

    palette = colors or CHART_PALETTE
    plot = chart.plots[0]
    for i, s in enumerate(plot.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = palette[i % len(palette)]

    _style_axis_fonts(chart, font_size=font_size)
    if show_value:
        plot.has_data_labels = True
        _style_data_labels(plot.data_labels, font_size=font_size,
                           number_format=number_format)

    legend = (len(series) > 1) if show_legend is None else show_legend
    if legend:
        _style_legend(chart, position=XL_LEGEND_POSITION.BOTTOM,
                      font_size=font_size)
    else:
        chart.has_legend = False
    return chart


def line_chart(slide, x, y, w, h, *, categories, series, colors=None,
               markers=True, show_value=False, show_legend=None,
               number_format=None, line_width=2.0, font_size=10):
    """折れ線グラフを配置する。

    series は (系列名, 値リスト) のリスト。colors は系列ごとの線色（RGBColor）。
    markers=True で各データ点にマーカーを出す。line_width は線幅（pt）。
    その他の引数は bar_chart と同様。
    """
    ct = XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE
    cd = _build_category_data(categories, series, number_format)
    chart = _add_chart(slide, ct, x, y, w, h, cd).chart

    palette = colors or CHART_PALETTE
    plot = chart.plots[0]
    for i, s in enumerate(plot.series):
        col = palette[i % len(palette)]
        s.format.line.color.rgb = col
        s.format.line.width = Pt(line_width)
        if markers:
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = col
            s.marker.format.line.color.rgb = col

    _style_axis_fonts(chart, font_size=font_size)
    if show_value:
        plot.has_data_labels = True
        _style_data_labels(plot.data_labels, font_size=font_size,
                           number_format=number_format)

    legend = (len(series) > 1) if show_legend is None else show_legend
    if legend:
        _style_legend(chart, position=XL_LEGEND_POSITION.BOTTOM,
                      font_size=font_size)
    else:
        chart.has_legend = False
    return chart


def pie_chart(slide, x, y, w, h, *, categories, values, colors=None,
              show_value=True, show_percentage=False, show_legend=True,
              number_format=None, font_size=10):
    """円グラフを配置する。

    Args:
        categories: 扇のラベルのリスト。
        values: 値のリスト（categories と同順・同数）。
        colors: 扇ごとの色（RGBColor）のリスト。省略時は CHART_PALETTE。
        show_value: 値ラベルを表示する。
        show_percentage: 百分率ラベルを表示する。
        show_legend: 凡例表示（既定 True、位置は右）。
        number_format: 値ラベルの表示書式。
        font_size: ラベル・凡例の基準フォントサイズ（pt）。

    Returns:
        Chart オブジェクト。
    """
    cd = _build_category_data(categories, [("", values)], number_format)
    chart = _add_chart(slide, XL_CHART_TYPE.PIE, x, y, w, h, cd).chart

    palette = colors or CHART_PALETTE
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = palette[i % len(palette)]

    if show_value or show_percentage:
        plot.has_data_labels = True
        dl = plot.data_labels
        _style_data_labels(dl, font_size=font_size, number_format=number_format)
        dl.show_value = show_value
        dl.show_percentage = show_percentage

    if show_legend:
        _style_legend(chart, position=XL_LEGEND_POSITION.RIGHT,
                      font_size=font_size)
    else:
        chart.has_legend = False
    return chart


# ============================================================================
# flat components（本文コンポーネント・フラット標準）
# ============================================================================

def card(slide, x, y, w, h, *, title=None, body=None, chip=None, kicker=None,
         accent=PRIMARY, fill=WHITE, line=BORDER, line_w=0.75, radius=0.045,
         pad=0.5, title_color=TEXT, body_color=TEXT_MUTED,
         title_size=T_H2, body_size=T_BODY, kicker_size=10.0,
         line_spacing=1.35, jp_font=JP_FONT):
    """フラットな角丸カード（影なし・ヘアライン枠）。

    色ベタの帯見出しを使わず、アクセント色は左上の小さなチップに限定し、見出しは
    濃色テキストで置くのが今風の要。``chip`` を渡すとアクセント色の角丸チップを
    左上に配置し、その右へ ``title``／``kicker``、下に ``body`` を流す。chip=None
    なら左上から title を開始する。返り値はカード本体シェイプ。

    Args:
        x, y, w, h: カード矩形（cm）。
        title: 見出し（濃色テキスト）。
        body: 本文（文字列または複数行リスト）。
        chip: アクセントチップ内の短いラベル（番号など）。None でチップなし。
        kicker: 見出し下の補足語（アクセント色・小さめ）。
        accent: チップ／kicker のアクセント色（RGBColor）。
    """
    card_sh = shape_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                        fill=fill, line=line, line_w=line_w)
    try:
        card_sh.adjustments[0] = radius
    except (IndexError, KeyError):
        pass
    register_container(slide, x, y, w, h)   # 子要素はみ出し検査の親として登録
    inner_w = w - 2 * pad
    ty = y + pad
    if chip is not None:
        chip_sz = 1.0
        c = shape_box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + pad, ty, chip_sz, chip_sz,
                      text=chip, fill=accent, color=WHITE, size=18, bold=True,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=JOSEFIN)
        try:
            c.adjustments[0] = 0.25
        except (IndexError, KeyError):
            pass
        tx = x + pad + chip_sz + 0.35
        head_w = w - pad - chip_sz - 0.6
        if title is not None:
            textbox(slide, tx, ty - 0.05, head_w, 0.72, title, size=title_size,
                    bold=True, color=title_color, anchor=MSO_ANCHOR.MIDDLE, jp_font=jp_font)
        if kicker is not None:
            textbox(slide, tx, ty + 0.62, head_w, 0.62, kicker, size=kicker_size,
                    color=accent, jp_font=jp_font)
        body_y = ty + chip_sz + 0.55
    else:
        if title is not None:
            textbox(slide, x + pad, ty, inner_w, 0.72, title, size=title_size,
                    bold=True, color=title_color, jp_font=jp_font)
        body_y = ty + (0.85 if title is not None else 0.0)
        if kicker is not None:
            textbox(slide, x + pad, body_y, inner_w, 0.62, kicker, size=kicker_size,
                    color=accent, jp_font=jp_font)
            body_y += 0.7
    if body is not None:
        textbox(slide, x + pad, body_y, inner_w, h - (body_y - y) - pad, body,
                size=body_size, color=body_color, line_spacing=line_spacing,
                anchor=MSO_ANCHOR.TOP, jp_font=jp_font)
    return card_sh


def callout(slide, x, y, w, h, text, *, accent=PRIMARY, fill=PRIMARY_LIGHT,
            color=TEXT, size=None, bold=True, radius=None, jp_font=JP_FONT):
    """キーテイクアウェイ用の淡色バー＋左アクセントバー。

    反転白文字の色ベタ帯ではなく、淡色面＋左の細いアクセントバー＋濃色テキストで
    結論を示すフラット表現。返り値は下地シェイプ。

    下地は **直角の矩形**（RECTANGLE）にする。角丸の淡色ピルの上に直角のアクセント
    バーを重ねると、左端上下で角丸コーナーとバーの直角がかみ合わず、バーと下地の間に
    三日月状の隙間が生じて崩れて見える（design-guide §3.5「角丸×直角の入れ子禁止」）。
    下地とアクセントバーを共に直角にして左端を面一（つらいち）にすることで解消する。
    ``radius`` 引数は後方互換のため残すが無視する。"""
    if size is None:
        size = T_BODY + 2
    bar = shape_box(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=fill)
    shape_box(slide, MSO_SHAPE.RECTANGLE, x, y, 0.14, h, fill=accent)
    textbox(slide, x + 0.55, y, w - 1.0, h, text, size=size, bold=bold,
            color=color, anchor=MSO_ANCHOR.MIDDLE, jp_font=jp_font)
    return bar


# ============================================================================
# slide-type configurators
# ============================================================================

def configure_body(slide, *, section_label, title,
                   label_color=TEXT, title_color=TEXT,
                   label_size=12, title_size=24):
    """Layout 2 'ページ番号+タイトル社外秘' のプレースホルダを設定。
    上部の小ラベル（y~0.94）に section_label、下のタイトル（y~1.42）に title。"""
    for ph in slide.placeholders:
        try:
            top_cm = Emu(ph.top).cm
        except (TypeError, ValueError):
            continue
        if 0.5 < top_cm < 1.2:
            fill_tf(ph.text_frame, section_label,
                    size=label_size, bold=False, color=label_color,
                    align=PP_ALIGN.LEFT)
            # lstStyle の hanging indent (marL/indent) をテンプレート通りにリセット
            for p in ph.text_frame.paragraphs:
                pPr = p._p.get_or_add_pPr()
                pPr.set('marL', '0')
                pPr.set('indent', '0')
            _register(slide, "ph_fixed", Emu(ph.left).cm, top_cm,
                      Emu(ph.width).cm, Emu(ph.height).cm, section_label,
                      size=label_size)
        elif 1.2 <= top_cm < 2.0:
            fill_tf(ph.text_frame, title,
                    size=title_size, bold=True, color=title_color,
                    align=PP_ALIGN.LEFT)
            for p in ph.text_frame.paragraphs:
                pPr = p._p.get_or_add_pPr()
                pPr.set('marL', '0')
                pPr.set('indent', '0')
            _register(slide, "ph_fixed", Emu(ph.left).cm, top_cm,
                      Emu(ph.width).cm, Emu(ph.height).cm, title,
                      size=title_size, bold=True)


def configure_chapter(slide, *, chapter_num, title, subtitle="",
                      num_size=25, title_size=20, subtitle_size=12):
    """Layout 1 'TITLE_AND_BODY' （章扉）のプレースホルダ設定。
    chapter_num 例: 'PART  02'。subtitle はメインタイトル下の小文字で表示。
    章番号は Josefin Sans（テンプレート定義フォント）で描画される。"""
    for ph in slide.placeholders:
        try:
            top_cm = Emu(ph.top).cm
        except (TypeError, ValueError):
            continue
        if 3.5 < top_cm < 4.5:
            fill_tf(ph.text_frame, chapter_num,
                    size=num_size, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, font=JOSEFIN)
            _register(slide, "ph_fixed", Emu(ph.left).cm, top_cm,
                      Emu(ph.width).cm, Emu(ph.height).cm, chapter_num,
                      size=num_size, bold=True, font=JOSEFIN)
        elif 8.0 < top_cm < 9.0:
            lines = [title]
            if subtitle:
                lines.append(subtitle)
            fill_tf(ph.text_frame, lines,
                    size=title_size, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, line_spacing=1.25)
            _register(slide, "ph_fixed", Emu(ph.left).cm, top_cm,
                      Emu(ph.width).cm, Emu(ph.height).cm, lines,
                      size=title_size, bold=True, line_spacing=1.25)
            if subtitle:
                p_sub = ph.text_frame.paragraphs[1]
                for r in p_sub.runs:
                    r.font.size = Pt(subtitle_size)
                    r.font.bold = False


def update_cover(slide, *, lines):
    """v3 サンプル表紙（テンプレのロゴ + 右側テキストボックス + CONFIDENTIAL）の
    右側テキストボックスを差し替え。

    lines: list[dict] 各要素は
        {"text": str, "size": int, "bold": bool, "color": RGBColor, "italic": bool}
    例:
        [
            {"text": "株式会社サンプル銀行 御中", "size": 12, "color": WHITE},
            {"text": "AIチャットボット セキュリティ構成のご提案",
             "size": 17, "bold": True, "color": WHITE},
            {"text": "外部公開環境において…", "size": 10, "italic": True, "color": WHITE},
            {"text": "2026年5月  ／  Ver. 1.0", "size": 10, "color": WHITE},
        ]
    """
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        txt = shp.text_frame.text
        # サンプル表紙の主要テキストボックス（"御中" or "ご提案資料" を含む）を識別
        if "御中" in txt or "ご提案資料" in txt:
            line_h_cm = sum(line.get("size", 14) * 1.4 / 72 * 2.54 for line in lines)
            new_h = Cm(line_h_cm + 0.6)
            shp.height = new_h
            shp.top = int((Cm(SLIDE_H_CM) - new_h) / 2)
            tf = shp.text_frame
            tf.word_wrap = True
            paras = list(tf.paragraphs)
            for p in paras[1:]:
                p._p.getparent().remove(p._p)
            p0 = paras[0]
            for r in list(p0.runs):
                r._r.getparent().remove(r._r)
            for i, line in enumerate(lines):
                p = p0 if i == 0 else tf.add_paragraph()
                run = p.add_run()
                _set_run(run, line["text"],
                         size=line.get("size", 14),
                         bold=line.get("bold", False),
                         color=line.get("color", WHITE),
                         italic=line.get("italic", False),
                         font=line.get("font", JOSEFIN))
            _register(slide, "ph_fixed", Emu(shp.left).cm, Emu(shp.top).cm,
                      Emu(shp.width).cm, Emu(shp.height).cm,
                      [ln["text"] for ln in lines],
                      size=max(ln.get("size", 14) for ln in lines),
                      line_spacing=1.4, autosize_v=True,
                      font=lines[0].get("font", JOSEFIN) if lines else JOSEFIN)
            return shp
    return None


def _body_layout_name(prs):
    """本文レイアウトの名前を prs の slide_layouts[L_BODY] から解決する。

    レイアウト名のリテラル照合はリブランド（雛形差し替え）で改名されると
    無音で判定が外れるため、REBRAND.md が契約とする index（L_BODY=2）から
    毎回動的に取得する。名前が何であれ index 構成を守った雛形なら追従する。"""
    return prs.slide_layouts[L_BODY].name


def _is_body_slide(slide):
    """slide が本文レイアウト（L_BODY）由来かを判定する。"""
    prs = slide.part.package.main_document_part.presentation
    return slide.slide_layout.name == _body_layout_name(prs)


def add_page_number(slide, num, *, body_text_color=WHITE, accent_text_color=PRIMARY):
    """右上のページ番号サークル位置 (23.77, 0.46) に番号テキストを挿入。

    Layout 2 (本文) は主色サークル → 既定 WHITE 文字。
    Layout 1 (章扉) は白サークル → 既定 PRIMARY 文字。
    body_text_color/accent_text_color で上書き可能。"""
    if _is_body_slide(slide):
        text_color = body_text_color
    else:
        text_color = accent_text_color
    tb = slide.shapes.add_textbox(Cm(23.77), Cm(0.46), Cm(1.16), Cm(1.06))
    tf = tb.text_frame
    tf.margin_left = Cm(0)
    tf.margin_right = Cm(0)
    tf.margin_top = Cm(0)
    tf.margin_bottom = Cm(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.fill.background()
    tb.line.fill.background()
    fill_tf(tf, str(num), size=23, bold=True, color=text_color,
            align=PP_ALIGN.CENTER, font=JOSEFIN)
    return tb


# ============================================================================
# icon helper  (Iconify API — https://iconify.design, 15万種以上のアイコン)
# ============================================================================

_icon_cache: dict = {}


def _icon_cache_dir():
    """アイコン PNG の永続キャッシュディレクトリ（ユーザースコープ）。

    プロセス再起動・別案件・別 worktree でも共有できるよう
    XDG 慣習に従い ~/.cache/ksr-slides/icons/ に置く。
    """
    from pathlib import Path
    d = Path.home() / ".cache" / "ksr-slides" / "icons"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _icon_png(icon_id: str, color: str, size_px: int = 256) -> bytes:
    """Iconify API から SVG を取得し PNG bytes に変換。

    キャッシュ階層:
      1. プロセス内メモリ（_icon_cache）
      2. ディスク（~/.cache/ksr-slides/icons/）— 再ラン・別案件でも共有
      3. HTTP 取得（Iconify API）— キャッシュミス時のみ
    """
    key = (icon_id, color, size_px)
    if key in _icon_cache:
        return _icon_cache[key]
    safe_id = icon_id.replace(":", "__").replace("/", "_")
    safe_color = color.lstrip("#")
    disk_path = _icon_cache_dir() / f"{safe_id}__{safe_color}__{size_px}.png"
    if disk_path.exists():
        png = disk_path.read_bytes()
        _icon_cache[key] = png
        return png
    try:
        import requests as _req
        import cairosvg as _svg
    except ImportError as exc:
        raise ImportError(
            "add_icon() には requests と cairosvg が必要です: "
            "pip install requests cairosvg"
        ) from exc
    prefix, name = icon_id.split(":")
    url = (
        f"https://api.iconify.design/{prefix}/{name}.svg"
        f"?color=%23{safe_color}&width={size_px}&height={size_px}"
    )
    resp = _req.get(url, timeout=10)
    resp.raise_for_status()
    png = _svg.svg2png(bytestring=resp.content, output_width=size_px, output_height=size_px)
    disk_path.write_bytes(png)
    _icon_cache[key] = png
    return png


def _hex(c) -> str:
    """RGBColor → "#RRGGBB" hex 文字列（add_icon の color= はブランド色定数でなく
    hex を取るため、brand.py の定数から変換して使う。RGBColor は str サブクラス）。"""
    return f"#{c}"


# ── 定義済みアイコン定数  (Iconify icon_id, hex_color) ─────────────────────
# 任意の Iconify アイコンを "セット名:アイコン名" 形式で追加可能。
# アイコン検索 → https://icon-sets.iconify.design/
# 既定色はブランド色定数から導出するため、brand.py を書き換えれば追従する。
ICON_TREE        = ("mdi:sitemap",          _hex(SECONDARY))  # ツリー/構成図
ICON_CERT        = ("mdi:certificate",      _hex(SECONDARY))  # 証明書/書類
ICON_EYE         = ("mdi:eye",              _hex(SECONDARY))  # 監視/監査
ICON_KEY         = ("mdi:key",              _hex(PRIMARY))    # 鍵/認証
ICON_LOCK        = ("mdi:lock",             _hex(PRIMARY))    # 錠前/パスワード
ICON_FINGERPRINT = ("mdi:fingerprint",      _hex(PRIMARY))    # 指紋/MFA
ICON_PHONE       = ("mdi:cellphone",        _hex(SECONDARY))  # スマートフォン/端末
ICON_USER_SHIELD = ("mdi:shield-account",   _hex(PRIMARY))    # アカウント保護/担当者
ICON_CHECK       = ("mdi:check-circle",     _hex(SUCCESS))    # OK/推奨（SUCCESS）
ICON_X           = ("mdi:close-circle",     _hex(DANGER))     # NG/問題点（DANGER）
ICON_ARROW       = ("mdi:arrow-right-bold", _hex(PRIMARY))    # フロー/遷移


def add_icon(slide, icon_spec, x: float, y: float, size: float = 1.0,
             color: str | None = None):
    """Iconify アイコンをスライドに配置する（cm 座標）。

    icon_spec: ICON_* 定数 (tuple) または "mdi:shield-lock" 形式の str。
               str を渡す場合は color 引数で色を指定（省略時: SECONDARY）。
    color: 省略すると ICON_* 定数の既定色を使用。"#EC6739" などの hex string。
    size: 正方形サイズ（cm）。
    """
    import io as _io
    if isinstance(icon_spec, tuple):
        icon_id, default_color = icon_spec
    else:
        icon_id, default_color = icon_spec, _hex(SECONDARY)
    png_bytes = _icon_png(icon_id, color or default_color)
    slide.shapes.add_picture(_io.BytesIO(png_bytes), Cm(x), Cm(y), Cm(size), Cm(size))


# ============================================================================
# convenience: open template / save
# ============================================================================

def load_template(path=None):
    """雛形を開いて Presentation を返す。path 省略時は templates/<brand.TEMPLATE>。"""
    if path is None:
        from pathlib import Path
        here = Path(__file__).resolve().parent
        path = here / "templates" / TEMPLATE
    _FIT_REGISTRY.clear()   # 各生成ランごとに検知レジストリをリセット
    _BAND_REGISTRY.clear()
    _GEOM_REGISTRY.clear()
    _ARROW_REGISTRY.clear()
    _CONTAINER_REGISTRY.clear()
    return Presentation(str(path))


def reset_to_cover_only(prs):
    """雛形のサンプルスライド（章扉/本文/裏表紙）を削除し、表紙(idx 0)のみ残す。
    update_cover() で表紙を書き換えてから残りページを add_slide で追加する想定。"""
    for idx in sorted(range(1, len(prs.slides)), reverse=True):
        delete_slide(prs, idx)


def add_back_cover(prs):
    """テンプレートの裏表紙（ロゴ込み）を prs の末尾に複製する。
    add_slide(L_BACK) では画像が消えるため、表紙と同様にテンプレートから clone する方式。"""
    import copy
    from pathlib import Path
    here = Path(__file__).resolve().parent
    tmpl = Presentation(str(here / "templates" / TEMPLATE))
    src = tmpl.slides[3]  # テンプレートの裏表紙は idx=3

    slide = prs.slides.add_slide(prs.slide_layouts[L_COVER])  # TITLEレイアウト

    # 新スライドの shape tree を空にする
    sp_tree = slide.shapes._spTree
    for el in list(sp_tree)[2:]:
        sp_tree.remove(el)

    # 画像リレーションを移植（prs 内に同一 blob があれば再利用して重複を避ける）
    rId_map = {}
    for rel in src.part.rels.values():
        if "image" not in rel.reltype:
            continue
        src_blob = rel.target_part.blob
        existing_part = None
        for existing_slide in prs.slides:
            for er in existing_slide.part.rels.values():
                if "image" in er.reltype and er.target_part.blob == src_blob:
                    existing_part = er.target_part
                    break
            if existing_part:
                break
        target_part = existing_part if existing_part else rel.target_part
        rId_map[rel.rId] = slide.part.relate_to(target_part, rel.reltype)

    # src の各 shape を複製（r:embed の rId を置換してから追加）
    for el in list(src.shapes._spTree)[2:]:
        xml_str = _etree.tostring(copy.deepcopy(el)).decode()
        for old_rId, new_rId in rId_map.items():
            xml_str = xml_str.replace(f'r:embed="{old_rId}"', f'r:embed="{new_rId}"')
        sp_tree.append(_etree.fromstring(xml_str))


def finalize_page_numbers(prs, *, skip_first=True):
    """全スライドにページ番号テキストを挿入。
    skip_first=True なら 1ページ目（表紙）はスキップ
    （サンプル表紙の placeholder '1' を再利用するため）。"""
    for i, s in enumerate(prs.slides):
        if skip_first and i == 0:
            continue
        add_page_number(s, i + 1)


# ============================================================================
# レイアウト検査  (枠はみ出し / 見切れ / ゾーン逸脱の機械検知)
#
#   validate_fit(prs)      … 生成時リント（PIL+Arial で折返しを近似計測）
#   validate_render(prs,p) … レンダ後突合（soffice→PDF を実測し意図文字列と照合）
#   validate(prs, path)    … 上記2つを実行し統一レポートを表示
# ============================================================================

SLACK_CM = 0.07           # 許容スラック（cm）
REL_SLACK = 0.04          # 許容スラック（相対）

# 「色帯」とみなす塗り矩形の最大高さ（cm）。これ以下の低い帯に textbox を
# 載せていて、その textbox が帯の下端を突き抜ける場合は確定的な視覚崩れ
# （色帯の外に文字が漏出）とみなし BAND_OVERFLOW=ERROR にする。背の高い
# 白カード（~8cm）はこの閾値を超えるため対象外＝従来どおり WARN のまま。
BAND_MAX_H = 3.0

# 1行の縦寸法を pt 単位で見積るときの係数。
# python-pptx / LibreOffice の Arial 標準は 1.2 だが、PowerPoint 実描画では
# CJK フォールバック（MS P ゴシック等）の ascent/descent が大きく、実効値は
# 1.3〜1.4 にぶれる。LibreOffice ベースの validate_render が PASS しても
# PowerPoint で枠を破る事案を防ぐため、保守側の 1.30 で見積る。
LINE_H_FACTOR = 1.30

# 個別 shape_box の親矩形に対する実測 bbox 逸脱の許容スラック（cm）
SHAPE_PAD = 0.05

# 矢じり付きコネクタ（矢印）の「軸の可視長」検査。矢じり（med）の長さは線幅に
# 比例し、実測で線幅(pt)の約 4 倍。総長から矢じり分を引いた軸の可視長が
# MIN_ARROW_SHAFT_CM 未満だと、矢じりだけが潰れて向きが読めない矢印になる
# （バンド間の隙間が狭すぎる典型。geom のゾーン検査では拾えなかった盲点）。
ARROW_HEAD_K = 4.0          # 矢じり長 ≒ ARROW_HEAD_K × 線幅(pt)
MIN_ARROW_SHAFT_CM = 0.12   # 軸の可視長の下限（cm）

# register_container で登録した背景枠からの子要素はみ出し許容スラック（cm）。
# 通常の子は枠端から 0.2cm 以上の padding を持つため、これより小さいはみ出し
# （枠線をまたぐ崩れ）を検知できるよう小さめにする。
CONTAINER_SLACK_CM = 0.03

_FONT_CACHE: dict = {}
_FONT_PATH_CACHE: dict = {}
_FONT_WARNED: set = set()
# 解決系が使えない環境（fontconfig 無し）向けの最終フォールバック。
_ARIAL_CANDIDATES = {
    True: ["/System/Library/Fonts/Supplemental/Arial Bold.ttf"],
    False: ["/System/Library/Fonts/Supplemental/Arial.ttf",
            "/Library/Fonts/Arial.ttf"],
}
# fontconfig 不在時に family 名でファイルを探す標準的なフォントディレクトリ。
_FONT_DIRS = [
    "/System/Library/Fonts",
    "/System/Library/Fonts/Supplemental",
    "/Library/Fonts",
    os.path.expanduser("~/Library/Fonts"),
    "/usr/share/fonts",
    "/usr/local/share/fonts",
    os.path.expanduser("~/.fonts"),
    r"C:\Windows\Fonts",
]


def _warn_font_once(family: str, msg: str) -> None:
    """フォント解決の縮退を1 family につき一度だけ stderr へ（沈黙させない）。"""
    if family in _FONT_WARNED:
        return
    _FONT_WARNED.add(family)
    import sys
    print(f"[slides._font] {family}: {msg}", file=sys.stderr)


def _family_matches(matched: str, requested: str) -> bool:
    """解決されたフォント family が要求 family と実質同一か（空白無視・先頭要素比較）。"""
    m = matched.split(",")[0].strip().lower()
    r = requested.split(",")[0].strip().lower()
    return m == r or m.replace(" ", "") == r.replace(" ", "")


def _fc_match(family: str, bold: bool):
    """fontconfig(fc-match) で family→(ファイル, 実 family) を解決。
    fontconfig は LibreOffice が Linux で使う解決系と同一なので、レンダラの
    フォント代替挙動に一致した計測ができる。無ければ (None, None)。"""
    import shutil
    fc = shutil.which("fc-match")
    if not fc:
        return None, None
    import subprocess
    try:
        pat = f"{family}:weight={'bold' if bold else 'regular'}"
        r = subprocess.run([fc, "-f", "%{file}\t%{family}", pat],
                           capture_output=True, text=True, timeout=5)
    except Exception:  # noqa: BLE001  解決失敗は握り潰さず None 返し（呼出側で警告）
        return None, None
    if r.returncode != 0 or not r.stdout.strip():
        return None, None
    path, _, fam = r.stdout.partition("\t")
    path, fam = path.strip(), fam.strip()
    return (path, fam) if path and os.path.exists(path) else (None, None)


def _glob_font_file(family: str, bold: bool):
    """fontconfig 不在時のフォールバック：フォントディレクトリを family 名で走査。"""
    import glob
    key = family.lower().replace(" ", "")
    if not key:
        return None
    cands = []
    for d in _FONT_DIRS:
        if not os.path.isdir(d):
            continue
        for ext in ("ttf", "ttc", "otf"):
            for p in glob.glob(os.path.join(d, f"*.{ext}")):
                stem = os.path.splitext(os.path.basename(p))[0].lower().replace(" ", "")
                if key in stem:
                    cands.append((p, "bold" in stem))
    if not cands:
        return None
    for p, is_bold in cands:  # 太字要求には bold ファイルを優先
        if is_bold == bool(bold):
            return p
    return cands[0][0]


def _resolve_font_path(family: str, bold: bool):
    """family(+bold)→(実ファイルパス, 実 family)。fontconfig→dir 走査→Arial の順。
    どれも見つからなければ (None, None)。結果は (family,bold) 単位でキャッシュ。"""
    family = family or FONT
    key = (family, bool(bold))
    if key in _FONT_PATH_CACHE:
        return _FONT_PATH_CACHE[key]
    path, matched = _fc_match(family, bold)
    if path is None:
        p = _glob_font_file(family, bold)
        if p:
            path, matched = p, family
    if path is None:  # 最終手段：同梱の Arial 候補
        for p in _ARIAL_CANDIDATES[bool(bold)] + _ARIAL_CANDIDATES[False]:
            if os.path.exists(p):
                path, matched = p, "Arial"
                break
    _FONT_PATH_CACHE[key] = (path, matched)
    return path, matched


def _font(family: str, size_pt: float, bold: bool):
    """PIL ImageFont を ×4 スケールで取得（サブpt精度のため）。
    family を実フォントファイルへ解決して計測に用いる（旧実装の Arial 固定を廃止）。
    - 解決 family が要求と違う（未インストールで代替された）→ WARN（レンダラも代替する）
    - どのファイルにも解決できない → WARN して PIL default（幅は不正確だが沈黙しない）"""
    family = family or FONT
    key = (family, round(size_pt, 1), bool(bold))
    f = _FONT_CACHE.get(key)
    if f is not None:
        return f
    from PIL import ImageFont
    path, matched = _resolve_font_path(family, bold)
    if path is None:
        _warn_font_once(family, "計測用フォントファイルを解決できず PIL default を使用"
                                "（幅推定は不正確。validate_render / 目視で確認を）")
        f = ImageFont.load_default()
    else:
        if matched and not _family_matches(matched, family):
            _warn_font_once(
                family,
                f"未インストール。'{matched}' の字幅で近似計測（レンダラも同様に代替）。"
                "固定幅枠は validate_render / 目視で要確認")
        f = ImageFont.truetype(path, max(4, int(round(size_pt * 4))))
    _FONT_CACHE[key] = f
    return f


def _is_cjk(ch: str) -> bool:
    """Unicode コードポイントから CJK 文字（ひらがな・カタカナ・漢字・全角等）か判定する。"""
    o = ord(ch)
    return (0x3040 <= o <= 0x30FF or 0x3400 <= o <= 0x9FFF
            or 0xAC00 <= o <= 0xD7A3 or 0xFF00 <= o <= 0xFFEF
            or 0x3000 <= o <= 0x303F)


def _latin_w(s: str, size_pt: float, bold: bool, family: str = None) -> float:
    """latin 文字列の幅（pt）。CJK は呼び出し側で 1em として扱う。
    family は計測に使う実フォント（未指定は既定 latin フォント FONT）。"""
    if not s:
        return 0.0
    f = _font(family or FONT, size_pt, bold)
    try:
        return f.getlength(s) / 4.0
    except Exception:  # 古い PIL
        bb = f.getbbox(s)
        return (bb[2] - bb[0]) / 4.0


def _wrap_lines(text: str, max_w_pt: float, size_pt: float, bold: bool,
                family: str = None):
    """LibreOffice 準拠の折返しを近似。
    latin は空白でのみ改行、CJK は1文字ごとに改行可。
    family は latin 幅計測に使う実フォント（未指定は既定 FONT）。
    returns (行リスト, 最大行幅pt)。"""
    family = family or FONT
    space_w = _latin_w(" ", size_pt, bold, family)
    words: list[str] = []
    buf = ""
    for ch in text:
        if ch in (" ", "\t", "　"):
            if buf:
                words.append(buf)
                buf = ""
            words.append(" ")
        elif _is_cjk(ch):
            if buf:
                words.append(buf)
                buf = ""
            words.append(ch)
        else:
            buf += ch
    if buf:
        words.append(buf)

    def ww(t: str) -> float:
        if t == " ":
            return space_w
        if len(t) == 1 and _is_cjk(t):
            return size_pt
        return _latin_w(t, size_pt, bold, family)

    lines: list[str] = []
    line, lw, maxw = "", 0.0, 0.0
    for tok in words:
        tw = ww(tok)
        if tok == " ":
            if line == "":
                continue
            if lw + tw <= max_w_pt:
                line += " "
                lw += tw
            else:
                lines.append(line)
                maxw = max(maxw, lw)
                line, lw = "", 0.0
            continue
        if line != "" and lw + tw > max_w_pt:
            lines.append(line)
            maxw = max(maxw, lw)
            line, lw = "", 0.0
        line += tok
        lw += tw
    if line != "" or not lines:
        lines.append(line)
        maxw = max(maxw, lw)
    return lines, maxw


class Violation:
    """レイアウト検査の検知結果 1 件を表すデータオブジェクト。

    Attributes:
        slide: 1-origin ページ番号。
        severity: 'ERROR'（完了不可）または 'WARN'（要確認）。
        code: 'OVERFLOW_V' / 'OVERFLOW_H' / 'ZONE' / 'CLIP' / 'LOST' /
              'PAGE_MISSING' / 'SHAPE_OVERFLOW' / 'WRAP' /
              'BAND_OVERFLOW' / 'ARROW_TOO_SHORT' / 'CONTAINER_OVERFLOW'。
        kind: 要素種別（'textbox' / 'shape' / 'cell' / 'ph_fixed' /
              'connector' / 'picture' / 'chart'）。
        rect: 要素の矩形 (x, y, w, h) cm タプル。
        detail: 検知の詳細説明文字列。
    """
    __slots__ = ("slide", "severity", "code", "kind", "rect", "detail")

    def __init__(self, slide, severity, code, kind, rect, detail):
        self.slide = slide          # 1-origin ページ番号
        self.severity = severity    # 'ERROR' / 'WARN'
        self.code = code            # 'OVERFLOW_V' / 'OVERFLOW_H' / 'ZONE' / 'CLIP' / 'LOST' / 'PAGE_MISSING' / 'SHAPE_OVERFLOW' / 'WRAP' / 'BAND_OVERFLOW'
        self.kind = kind
        self.rect = rect            # (x, y, w, h) cm
        self.detail = detail

    def __repr__(self):
        x, y, w, h = self.rect
        return (f"P{self.slide:>2} [{self.severity}] {self.code:<10} "
                f"{self.kind:<8} @({x:.1f},{y:.1f} {w:.1f}x{h:.1f}) "
                f"{self.detail}")


def _slide_index_map(prs):
    """slide_id → 0-origin インデックスの辞書を返す（_FIT_REGISTRY 照合用）。"""
    return {s.slide_id: i for i, s in enumerate(prs.slides)}


def _snip(lines, n=40):
    """複数行テキストを n 文字で切り詰めて検査レポート用の表示文字列を返す。"""
    t = " / ".join(lines)
    return t if len(t) <= n else t[:n] + "…"


def _find_overflow_band(slide_id, x, y, w, need_h_cm):
    """textbox（左上 x,y / 幅 w / 自動伸長後の必要高 need_h_cm）が、内包する
    低い「色帯」の下端を突き抜けているかを調べ、該当する帯の dict を返す（無ければ None）。

    色帯の外に文字が漏れるのは目視確認に頼らず確定できる視覚崩れなので、
    これに該当する縦あふれは WARN ではなく ERROR(BAND_OVERFLOW) に格上げする。
    背の高い白カード（h>BAND_MAX_H）は帯とみなさないため、カード内で保守的な
    need_h が枠を超えるだけの空振り（例: P3 型）は従来どおり WARN のまま。
    """
    eps = 0.10
    hit = None
    for b in _BAND_REGISTRY:
        if b["slide_id"] != slide_id:
            continue
        if b["h"] > BAND_MAX_H:                      # 背の高いカードは帯ではない
            continue
        if not (b["x"] - eps <= x and x + w <= b["x"] + b["w"] + eps):
            continue                                  # 帯の幅に収まっていない
        if not (b["y"] - eps <= y < b["y"] + b["h"]):
            continue                                  # textbox が帯内から始まっていない
        band_bottom = b["y"] + b["h"]
        if y + need_h_cm <= band_bottom + SLACK_CM:
            continue                                  # 帯の下端を越えていない
        # 複数の帯に該当する場合は最も内側（下端が高い）の帯を採用
        if hit is None or band_bottom < hit["y"] + hit["h"]:
            hit = b
    return hit


def validate_fit(prs):
    """生成時リント（レンダリング不要）。レジストリの各要素について
    PIL＋その要素の実フォント（rec["font"]）で折返しを近似計測し、
    縦あふれ/横はみ出し/ゾーン逸脱を検出。"""
    idx_map = _slide_index_map(prs)
    body_layout = _body_layout_name(prs)
    out: list[Violation] = []
    for rec in _FIT_REGISTRY:
        si = idx_map.get(rec["slide_id"])
        if si is None:
            continue
        page = si + 1
        size, bold, ls = rec["size"], rec["bold"], rec["ls"]
        family = rec.get("font", FONT)
        kind = rec["kind"]
        x, y, w, h = rec["x"], rec["y"], rec["w"], rec["h"]
        eff_w_pt = max(1.0, (w - rec["ml"] - rec["mr"]) * PT_PER_CM)

        total_lines, max_line_pt = 0, 0.0
        for para in rec["lines"]:
            ls_, mw = _wrap_lines(para, eff_w_pt, size, bold, family)
            total_lines += len(ls_)
            max_line_pt = max(max_line_pt, mw)

        line_h_pt = size * LINE_H_FACTOR * ls
        need_h_cm = (total_lines * line_h_pt) / PT_PER_CM + rec["mt"] + rec["mb"]
        slack_v = max(SLACK_CM, h * REL_SLACK)

        # 横はみ出し（折返し後の最長行が枠内幅を超える＝分割不能な語/字）。
        # 細幅シェイプ（pill/badge: w<3cm）はスラックを極小にする。
        # 通常スラック（rel 4% + abs 0.07cm）だと 1.45cm 枠 vs 1.47cm
        # テキストのような hairline 超過を見逃すが、実描画では確実に折返す。
        if kind == "shape" and w < 3.0:
            h_rel_slack = 0.0
            h_abs_slack_pt = 0.6   # ≒0.02cm
        else:
            h_rel_slack = REL_SLACK
            h_abs_slack_pt = 2.0   # ≒0.07cm
        if max_line_pt > eff_w_pt * (1 + h_rel_slack) + h_abs_slack_pt:
            over = max_line_pt / PT_PER_CM - (w - rec["ml"] - rec["mr"])
            out.append(Violation(
                page, "ERROR", "OVERFLOW_H", kind, (x, y, w, h),
                f"横はみ出し +{over:.2f}cm 「{_snip(rec['lines'])}」"))

        # 縦あふれ／見切れ（autosize_v はビルダーが枠高を自動調整するため対象外）
        if not rec.get("autosize_v") and need_h_cm > h + slack_v:
            if kind == "shape" or kind == "cell":
                # 固定矩形/表セルは確実にクリップする → ERROR
                out.append(Violation(
                    page, "ERROR", "OVERFLOW_V", kind, (x, y, w, h),
                    f"縦あふれ/見切れ 必要{need_h_cm:.2f}>枠{h:.2f}cm "
                    f"「{_snip(rec['lines'])}」"))
            elif kind == "ph_fixed":
                # テンプレ固定枠は PPT 側で自動縮小されうる → 一次は WARN
                # （真の見切れは validate_render の CLIP=ERROR で確定検知）
                out.append(Violation(
                    page, "WARN", "OVERFLOW_V", kind, (x, y, w, h),
                    f"固定枠に対し文字量過多の可能性 必要{need_h_cm:.2f}>"
                    f"枠{h:.2f}cm（テキスト短縮を検討）「{_snip(rec['lines'])}」"))
            else:  # textbox は縦に自動伸長
                band = _find_overflow_band(rec["slide_id"], x, y, w, need_h_cm)
                if band is not None:
                    # 色帯の下端を突き抜けて文字が漏出 → 確定的な崩れ → ERROR
                    leak = (y + need_h_cm) - (band["y"] + band["h"])
                    out.append(Violation(
                        page, "ERROR", "BAND_OVERFLOW", kind, (x, y, w, h),
                        f"色帯({band['h']:.2f}cm)からテキストが下端へ{leak:.2f}cm漏出 "
                        f"必要{need_h_cm:.2f}>枠{h:.2f}cm「{_snip(rec['lines'])}」"))
                else:  # 帯外（カード内等）の縦伸長 → 隣接重なりの WARN
                    out.append(Violation(
                        page, "WARN", "OVERFLOW_V", kind, (x, y, w, h),
                        f"縦に伸長 必要{need_h_cm:.2f}>枠{h:.2f}cm "
                        f"（隣接要素と重なる恐れ）「{_snip(rec['lines'])}」"))

        # コンテンツゾーン逸脱（本文レイアウトの自由配置要素のみ）
        if kind in ("textbox", "shape", "cell") and \
                rec["layout"] == body_layout:
            bottom = y + (need_h_cm if kind == "textbox" else h)
            msgs = []
            if x < X_MIN - SLACK_CM:
                msgs.append(f"x={x:.2f}<{X_MIN}")
            if x + w > X_MAX + SLACK_CM:
                msgs.append(f"右端={x + w:.2f}>{X_MAX}")
            if y < Y_MIN - SLACK_CM:
                msgs.append(f"y={y:.2f}<{Y_MIN}")
            if bottom > Y_MAX + SLACK_CM:
                msgs.append(f"下端={bottom:.2f}>{Y_MAX}")
            if msgs:
                out.append(Violation(
                    page, "ERROR", "ZONE", kind, (x, y, w, h),
                    "ゾーン逸脱 " + " ".join(msgs)
                    + f" 「{_snip(rec['lines'])}」"))

    # 非テキスト幾何要素（コネクタ/画像/グラフ）のコンテンツゾーン逸脱。
    # テキストを持たないため _FIT_REGISTRY ではなく _GEOM_REGISTRY を走査する。
    # 内部テキスト（軸・データラベル等）は検査しないが、要素の配置矩形が本文
    # ゾーンを外れていれば確定的な崩れなので ERROR(ZONE) とする。
    for grec in _GEOM_REGISTRY:
        si = idx_map.get(grec["slide_id"])
        if si is None or grec["layout"] != body_layout:
            continue
        x, y, w, h = grec["x"], grec["y"], grec["w"], grec["h"]
        msgs = []
        if x < X_MIN - SLACK_CM:
            msgs.append(f"x={x:.2f}<{X_MIN}")
        if x + w > X_MAX + SLACK_CM:
            msgs.append(f"右端={x + w:.2f}>{X_MAX}")
        if y < Y_MIN - SLACK_CM:
            msgs.append(f"y={y:.2f}<{Y_MIN}")
        if y + h > Y_MAX + SLACK_CM:
            msgs.append(f"下端={y + h:.2f}>{Y_MAX}")
        if msgs:
            out.append(Violation(
                si + 1, "ERROR", "ZONE", grec["kind"], (x, y, w, h),
                "ゾーン逸脱(非テキスト要素) " + " ".join(msgs)))

    # 矢じり付きコネクタ（矢印）の軸長検査。総長から矢じり分を引いた「軸の
    # 可視長」が不足すると、矢じりだけが潰れて向きが読めない矢印になる（geom の
    # ゾーン検査では拾えない）。データ欠落ではなく見栄えの崩れなので WARN で報告。
    for arec in _ARROW_REGISTRY:
        si = idx_map.get(arec["slide_id"])
        if si is None:
            continue
        head_cm = ARROW_HEAD_K * arec["width"] / PT_PER_CM
        shaft = arec["length"] - arec["heads"] * head_cm
        if shaft < MIN_ARROW_SHAFT_CM:
            x, y, w, h = arec["x"], arec["y"], arec["w"], arec["h"]
            out.append(Violation(
                si + 1, "WARN", "ARROW_TOO_SHORT", "connector", (x, y, w, h),
                f"矢印の軸が短く矢じりが潰れる 可視軸{shaft:.2f}<"
                f"{MIN_ARROW_SHAFT_CM}cm（総長{arec['length']:.2f}cm/"
                f"線幅{arec['width']:.1f}pt/矢じり{arec['heads']}）"
                f"→ 間隔を広げる/線を短くしない"))

    # コンテナ（カード/パネル）からの子要素はみ出し検査。register_container で
    # 登録した背景枠ごとに、「中心がその枠内にある・枠より小さい」子要素
    # （テキスト/チップ/図形）を集め、子の bbox が枠を越えていれば WARN。
    # geom のゾーン検査（ページ枠）では拾えない「枠内のはみ出し」を補完する。
    children = [(r["slide_id"], r["kind"], r["x"], r["y"], r["w"], r["h"],
                 _snip(r["lines"])) for r in _FIT_REGISTRY]
    children += [(g["slide_id"], g["kind"], g["x"], g["y"], g["w"], g["h"],
                  g["kind"]) for g in _GEOM_REGISTRY]
    for c in _CONTAINER_REGISTRY:
        si = idx_map.get(c["slide_id"])
        if si is None:
            continue
        cx0, cy0, cw, chh = c["x"], c["y"], c["w"], c["h"]
        cx1, cy1, c_area = cx0 + cw, cy0 + chh, cw * chh
        for sid, kind, ex, ey, ew, eh, label in children:
            if sid != c["slide_id"] or ew * eh >= c_area:
                continue   # 別スライド／枠と同等以上の大きさは子とみなさない
            ecx, ecy = ex + ew / 2, ey + eh / 2
            if not (cx0 <= ecx <= cx1 and cy0 <= ecy <= cy1):
                continue   # 中心が枠内にない＝この枠の子ではない
            msgs = []
            if ex < cx0 - CONTAINER_SLACK_CM:
                msgs.append(f"左{cx0 - ex:.2f}")
            if ex + ew > cx1 + CONTAINER_SLACK_CM:
                msgs.append(f"右{ex + ew - cx1:.2f}")
            if ey < cy0 - CONTAINER_SLACK_CM:
                msgs.append(f"上{cy0 - ey:.2f}")
            if ey + eh > cy1 + CONTAINER_SLACK_CM:
                msgs.append(f"下{ey + eh - cy1:.2f}")
            if msgs:
                out.append(Violation(
                    si + 1, "WARN", "CONTAINER_OVERFLOW", kind, (ex, ey, ew, eh),
                    f"枠({cx0:.1f},{cy0:.1f} {cw:.1f}x{chh:.1f})から子要素が"
                    f"はみ出し [{' '.join(msgs)}]cm 「{label}」"))
    return out


def _soffice_to_pdf(pptx_path, *, reuse=True):
    """LibreOffice (soffice) で .pptx を PDF に変換し、出力ファイルパスを返す。

    reuse=True（既定）のとき、既存 PDF の mtime が pptx 以上ならスキップして
    既存パスを返す。validate_render() と 5-b 視覚確認の二重変換を避けるため。
    prs.save() が必ず mtime を更新するため、stale を返すケースは実質ない。
    """
    import shutil
    import subprocess
    from pathlib import Path
    soffice = shutil.which("soffice") or "/opt/homebrew/bin/soffice"
    repo = Path(__file__).resolve().parent
    tmp = repo / ".claude" / "tmp" / "ksr-slides"
    tmp.mkdir(parents=True, exist_ok=True)
    pptx_path = Path(pptx_path)
    pdf = tmp / (pptx_path.stem + ".pdf")
    if reuse and pdf.exists() \
            and pdf.stat().st_mtime >= pptx_path.stat().st_mtime:
        return pdf
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(tmp), str(pptx_path)],
        check=True, capture_output=True)
    if not pdf.exists():
        raise RuntimeError(f"PDF 変換に失敗: {pdf}")
    return pdf


def render_pngs(pptx_path, *, dpi=120, out_dir=None):
    """pptx を PNG 群にレンダリングして保存パスのリストを返す（5-b 視覚確認用）。

    validate(prs, OUT, render=True) が既に作った PDF を再利用し、
    PyMuPDF で PNG 化する。soffice の二重起動と pdf2image (poppler) 依存を
    排除して 5-b を高速化するためのヘルパー。

    Args:
        pptx_path: 対象 .pptx のパス。
        dpi: PNG の解像度。既定 120（デザイン批評の十分な解像度）。
        out_dir: 出力ディレクトリ。省略時は PDF と同じ
                 .claude/tmp/ksr-slides/ に slide_NN.png として保存。

    Returns:
        list[Path]: 1ページ目→最終ページの順に PNG パス。
    """
    try:
        import fitz
    except ImportError as exc:
        raise ImportError(
            "render_pngs() には PyMuPDF が必要です: "
            "pip install pymupdf") from exc
    from pathlib import Path
    pdf_path = _soffice_to_pdf(pptx_path)
    out_dir = Path(out_dir) if out_dir else pdf_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    # soffice 生成 PDF の structure tree 軽微不整合に対する mupdf 警告は
    # レンダリング自体に影響しないので抑制してログを汚さない。
    fitz.TOOLS.mupdf_display_errors(False)
    doc = fitz.open(str(pdf_path))
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pngs = []
    for i, page in enumerate(doc):
        png = out_dir / f"slide_{i + 1:02d}.png"
        page.get_pixmap(matrix=mat).save(str(png))
        pngs.append(png)
    doc.close()
    return pngs


def _norm(s: str) -> str:
    """空白文字を全て除去して正規化する（レンダリング結果と意図テキストの照合用）。"""
    import re
    return re.sub(r"\s+", "", s)


def validate_render(prs, pptx_path):
    """レンダ後突合（実レンダラ準拠）。soffice→PDF を PyMuPDF で実測し、
    レジストリの意図テキストと照合。描画文字列が意図より短い箇所＝見切れ確定。"""
    try:
        import fitz
    except ImportError as exc:
        raise ImportError(
            "validate_render() には PyMuPDF が必要です: "
            "pip install pymupdf") from exc

    pdf_path = _soffice_to_pdf(pptx_path)
    # soffice 生成 PDF の structure tree 軽微不整合に対する mupdf 警告は
    # 検査結果に影響しないので抑制してログを汚さない。
    fitz.TOOLS.mupdf_display_errors(False)
    doc = fitz.open(str(pdf_path))
    slide_ids = [s.slide_id for s in prs.slides]
    body_layout = _body_layout_name(prs)
    out: list[Violation] = []

    # ページごとに span を (text, x0,y0,x1,y1) cm で収集
    page_spans = []
    for pg in doc:
        spans = []
        for blk in pg.get_text("dict")["blocks"]:
            for ln in blk.get("lines", []):
                for sp in ln["spans"]:
                    bx = [c / PT_PER_CM for c in sp["bbox"]]
                    spans.append((sp["text"], *bx))
        page_spans.append(spans)

    # PDF のページ欠落（スライド数より少ない）は、そのページ上の全要素の
    # 照合が成立しない確定異常。下のループの存在チェック（continue）だけだと
    # 黙殺されるため、ここで明示的に ERROR(PAGE_MISSING) にする。
    for missing_i in range(len(page_spans), len(slide_ids)):
        out.append(Violation(
            missing_i + 1, "ERROR", "PAGE_MISSING", "page",
            (0.0, 0.0, 0.0, 0.0),
            f"レンダ PDF にページが無い"
            f"（PDF {len(page_spans)}p ＜ スライド {len(slide_ids)}p）"))

    for rec in _FIT_REGISTRY:
        try:
            page_i = slide_ids.index(rec["slide_id"])
        except ValueError:
            continue
        if page_i >= len(page_spans):
            continue
        page = page_i + 1
        x, y, w, h = rec["x"], rec["y"], rec["w"], rec["h"]
        pad = 0.3
        cx0, cy0, cx1, cy1 = x - pad, y - pad, x + w + pad, y + h + pad
        matched = []
        for txt, sx0, sy0, sx1, sy1 in page_spans[page_i]:
            mx, my = (sx0 + sx1) / 2, (sy0 + sy1) / 2
            if cx0 <= mx <= cx1 and cy0 <= my <= cy1:
                matched.append((sy0, sx0, txt, sx0, sy0, sx1, sy1))
        matched.sort(key=lambda t: (round(t[0], 1), t[1]))
        rendered = _norm("".join(m[2] for m in matched))
        intended = _norm("".join(rec["lines"]))
        if not intended:
            continue

        if rendered == intended:
            pass
        elif not rendered:
            # 1 文字も描画されていない＝完全消失。一部見切れ（CLIP）より
            # 重い崩れなのに、従来は下の rendered 前提の分岐に入らず
            # 無言 PASS していた（shape_box の溢れ文字は LibreOffice が
            # 描画ごと消すため、まさにこのケースになる）。
            out.append(Violation(
                page, "ERROR", "LOST", rec["kind"], (x, y, w, h),
                f"完全消失: 意図テキストが 1 文字も描画されていない"
                f"（意図{len(intended)}字）「{_snip(rec['lines'])}」"))
        elif rendered and intended.startswith(rendered) \
                and len(rendered) < len(intended):
            lost = intended[len(rendered):]
            out.append(Violation(
                page, "ERROR", "CLIP", rec["kind"], (x, y, w, h),
                f"見切れ確定: 末尾「{lost[:24]}」が描画されていない "
                f"（意図{len(intended)}字→描画{len(rendered)}字）"))
        elif rendered and len(rendered) < len(intended) * 0.92 \
                and intended.startswith(rendered[:max(4, len(rendered) // 2)]):
            out.append(Violation(
                page, "WARN", "CLIP", rec["kind"], (x, y, w, h),
                f"描画文字列が意図より短い（要確認）"
                f"意図{len(intended)}字→描画{len(rendered)}字 "
                f"「{_snip(rec['lines'])}」"))

        # 実測ゾーン逸脱（本文レイアウトの自由配置要素のみ）
        if rec["kind"] in ("textbox", "shape", "cell") and \
                rec["layout"] == body_layout and matched:
            for _, _, txt, sx0, sy0, sx1, sy1 in matched:
                if not txt.strip():
                    continue
                if (sx0 < X_MIN - SLACK_CM or sx1 > X_MAX + SLACK_CM
                        or sy0 < Y_MIN - SLACK_CM or sy1 > Y_MAX + SLACK_CM):
                    out.append(Violation(
                        page, "ERROR", "ZONE", rec["kind"], (x, y, w, h),
                        f"実測ゾーン逸脱 描画bbox"
                        f"({sx0:.2f},{sy0:.2f})-({sx1:.2f},{sy1:.2f}) "
                        f"「{txt[:24]}」"))
                    break

        # 親 shape_box 矩形からの実測 bbox 逸脱
        #   PowerPoint / Google スライドで Arial メトリクスが LibreOffice
        #   より大きく出ることに起因する「LibreOffice では収まったように
        #   見えるが PowerPoint で枠を破る」事案を、LibreOffice 実測でも
        #   「枠端ギリギリ」として近似検知するための追加ゲート。
        #   textbox は縦に自動伸長する設計なので対象外、shape/cell
        #   （クリップ前提の固定枠）のみを検査する。
        #   matched は pad 0.3cm の検索範囲で取っているため隣接シェイプの
        #   span を拾うことがあり、ここでは「意図テキストの一部として描画
        #   された span」だけに絞り込んでから親矩形と照合する。
        if rec["kind"] in ("shape", "cell") and matched and intended:
            # 所属判定の閾値：matched は pad 0.3cm の検索範囲で取っている
            # ため隣接シェイプの span を含む。SHAPE_OVERFLOW では
            #   (a) 正規化テキストが意図テキストの部分文字列
            #   (b) かつ span 中心が親矩形 ± 0.10cm に収まる
            # の両方を満たす span のみを対象にする。
            own_pad = 0.10
            for _, _, txt, sx0, sy0, sx1, sy1 in matched:
                if not txt.strip():
                    continue
                nt = _norm(txt)
                if not nt or len(nt) < 3 or nt not in intended:
                    continue
                mx, my = (sx0 + sx1) / 2, (sy0 + sy1) / 2
                if not (x - own_pad <= mx <= x + w + own_pad
                        and y - own_pad <= my <= y + h + own_pad):
                    continue
                msgs = []
                if sx0 < x - SHAPE_PAD:
                    msgs.append(f"左 {sx0:.2f}<{x:.2f}")
                if sx1 > x + w + SHAPE_PAD:
                    msgs.append(f"右 {sx1:.2f}>{x + w:.2f}")
                if sy0 < y - SHAPE_PAD:
                    msgs.append(f"上 {sy0:.2f}<{y:.2f}")
                if sy1 > y + h + SHAPE_PAD:
                    msgs.append(f"下 {sy1:.2f}>{y + h:.2f}")
                if msgs:
                    out.append(Violation(
                        page, "ERROR", "SHAPE_OVERFLOW", rec["kind"],
                        (x, y, w, h),
                        "親矩形からの bbox 逸脱 " + " ".join(msgs)
                        + f" 「{txt[:24]}」"))
                    break

        # 1行想定の固定枠が複数 y-行に描画されていたら ERROR（WRAP）。
        #   pill/badge/タイルラベル等、validate_fit の OVERFLOW_H スラック内に
        #   収まる hairline 超過が実描画で折返している場合の最終ゲート。
        #   判定条件:
        #     (a) 1段落渡し（len(lines)==1）— list 複数渡しは明示的多行で対象外
        #     (b) 枠の有効高が 1 行分しかない（line_h × 1.8 未満）
        #         — 大きな本文ボックスへ 1 文字列を流し込んで自然折返しさせる
        #           本文ブロックを誤検知しないため
        if rec["kind"] in ("shape", "cell") and len(rec["lines"]) == 1 \
                and matched and intended:
            line_h_cm = (rec["size"] * LINE_H_FACTOR * rec["ls"]) / PT_PER_CM
            box_h_inner = max(0.0, h - rec["mt"] - rec["mb"])
            is_single_line_box = (line_h_cm > 0
                                  and box_h_inner < line_h_cm * 1.8)
            if is_single_line_box:
                own_pad = 0.10
                own_y_rows: list[float] = []
                for _, _, txt, sx0, sy0, sx1, sy1 in matched:
                    if not txt.strip():
                        continue
                    nt = _norm(txt)
                    if not nt or nt not in intended:
                        continue
                    mx, my = (sx0 + sx1) / 2, (sy0 + sy1) / 2
                    if not (x - own_pad <= mx <= x + w + own_pad
                            and y - own_pad <= my <= y + h + own_pad):
                        continue
                    # y-行へクラスタリング（0.15cm 以内なら同一行）
                    attached = False
                    for i, row_y in enumerate(own_y_rows):
                        if abs(my - row_y) <= 0.15:
                            own_y_rows[i] = (row_y + my) / 2
                            attached = True
                            break
                    if not attached:
                        own_y_rows.append(my)
                if len(own_y_rows) > 1:
                    out.append(Violation(
                        page, "ERROR", "WRAP", rec["kind"], (x, y, w, h),
                        f"1行想定の固定枠が{len(own_y_rows)}行に折返し "
                        f"「{rec['lines'][0][:24]}」"))
    doc.close()
    return out


def print_report(violations, *, title="レイアウト検査"):
    """検知結果を整形表示し、(error数, warn数) を返す。"""
    errs = [v for v in violations if v.severity == "ERROR"]
    warns = [v for v in violations if v.severity == "WARN"]
    print(f"\n===== {title} =====")
    if not violations:
        print("✅ PASS — 枠はみ出し・見切れ・ゾーン逸脱・親矩形逸脱は"
              "検出されませんでした。")
        return 0, 0
    for v in sorted(violations, key=lambda v: (v.slide, v.severity != "ERROR")):
        mark = "❌" if v.severity == "ERROR" else "⚠️"
        print(f"{mark} {v!r}")
    print(f"---- ERROR {len(errs)} 件 / WARN {len(warns)} 件 ----")
    if errs:
        print("❌ ERROR を解消するまで完了としないこと"
              "（座標・サイズ・テキスト量を調整して再生成）。")
    return len(errs), len(warns)


def validate(prs, pptx_path=None, *, render=True, strict=True):
    """生成時リント＋（任意で）レンダ後突合を実行し統一レポートを表示。

    pptx_path: render=True のとき soffice 変換対象（保存済み .pptx パス）。
    strict: 既定 True。ERROR があればレポート表示後に SystemExit(1) し、
            シェル・CI からも「必須ゲート」の失敗として機械的に見える。
            レポートだけ出して処理を続けたい場合のみ strict=False を指定。
    returns Violation のリスト。"""
    import sys
    violations = list(validate_fit(prs))
    if render:
        if pptx_path is None:
            raise ValueError("render=True には保存済み pptx_path が必要です")
        try:
            violations += validate_render(prs, pptx_path)
        except ImportError as exc:
            print(f"[validate] レンダ後突合をスキップ: {exc}", file=sys.stderr)
    n_err, _ = print_report(violations)
    if strict and n_err:
        sys.exit(1)
    return violations
